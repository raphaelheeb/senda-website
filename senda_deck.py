from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation()

def add_slide(title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    textbox = slide.shapes.add_textbox(Inches(1.2), Inches(2.2), Inches(7.5), Inches(2))
    tf = textbox.text_frame
    tf.clear()

    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.color.rgb = RGBColor(30, 30, 30)

    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(18)
        p2.font.color.rgb = RGBColor(100, 100, 100)

slides = [
("SENDA", "Swiss Performance Road & Gravel Bikes"),
("Engineering calm performance", "for a new generation of riders"),
("A refined approach to performance", "Built in Switzerland. Designed for Europe."),
("Calm • Precision • Performance • Timelessness", None),
("Born in Switzerland", "Precision, durability, design discipline"),
("The market has become visually loud", "Overbranding • Aggressive design"),
("The opportunity", "Minimal, design-led performance"),
("Positioning", "High performance meets refined design"),
("Product Vision", "One platform. Refined to its essence"),
("Platform Architecture", "Road • All-Road • Gravel"),
("Design Language", "Minimal • Precise • Timeless"),
("Product Experience", "Balanced ride feel and control"),
("Target Customer", "Design-conscious premium rider"),
("Market Context", "Premium growth and aesthetic demand rising"),
("Competitive Gap", "Space between mass brands and niche builders"),
("Pricing Strategy", "Premium, curated, high value"),
("Unit Economics", "Designed for sustainable margins"),
("Supply Chain", "Precision manufacturing partnerships"),
("Go-To-Market", "Brand • D2C • Retail • Community"),
("Switzerland First", "Strong premium entry market"),
("European Expansion", "Controlled rollout across key markets"),
("Channel Strategy", "D2C-first with selective retail"),
("Brand & Content", "Storytelling and design-led marketing"),
("Roadmap", "Launch → Expand → Scale")
]

for t, s in slides:
    add_slide(t, s)

prs.save("SENDA_Presentation_Final.pptx")

print("Presentation created successfully.")
